from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ──────────────────────────────────
BLUE       = RGBColor(0x00, 0x52, 0xCC)
BLUE_LIGHT = RGBColor(0xDE, 0xEB, 0xFF)
DARK       = RGBColor(0x17, 0x26, 0x48)
MUTED      = RGBColor(0x6B, 0x77, 0x8C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG         = RGBColor(0xF4, 0xF5, 0xF7)
GREEN      = RGBColor(0x00, 0x87, 0x5A)
AMBER      = RGBColor(0xFF, 0x8B, 0x00)
PURPLE     = RGBColor(0x64, 0x37, 0xD7)

CAPTURE = r"C:\Users\sam.lee\Desktop\잡일\WOW\Capture"
IMG = {
    'external': f"{CAPTURE}\\2026-03-27 10 41 19.png",
    'board':    f"{CAPTURE}\\2026-03-27 10 41 19 (2).png",
    'status':   f"{CAPTURE}\\2026-03-27 10 41 19 (3).png",
    'manage':   f"{CAPTURE}\\2026-03-27 10 41 19 (4).png",
    'drag':     f"{CAPTURE}\\2026-03-27 10 41 19 (5).png",
}

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=13, color=DARK, bullet='•'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def header_bar(slide, title, subtitle=None):
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, W, bar_h, BLUE)
    add_text(slide, title,
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.65),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.68), Inches(10), Inches(0.35),
                 font_size=14, color=RGBColor(0xB3, 0xD4, 0xFF))

def section_label(slide, text, left, top, color=BLUE):
    add_rect(slide, left, top, Inches(0.04), Inches(0.3), color)
    add_text(slide, text,
             left + Inches(0.12), top - Inches(0.02), Inches(5), Inches(0.35),
             font_size=14, bold=True, color=color)

def add_screenshot(slide, img_path, left, top, width, height, shadow=True):
    """스크린샷 삽입 + 흰 테두리 카드 효과"""
    # 흰 카드 배경
    pad = Inches(0.08)
    card = slide.shapes.add_shape(1,
        left - pad, top - pad, width + pad*2, height + pad*2)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDE, 0xEB, 0xFF)
    card.line.width = Pt(1)
    # 이미지
    slide.shapes.add_picture(img_path, left, top, width, height)

# ════════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BLUE)

add_rect(sl, 0, 0, W, H, BLUE)
# 반투명 장식 원
for (x, y, s, a) in [(Inches(10.5), Inches(0.5), Inches(4), 30),
                      (Inches(-1),   Inches(5),   Inches(3.5), 20)]:
    c = slide_shapes = sl.shapes.add_shape(9, x, y, s, s)  # oval
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(0x00, 0x3E, 0x9E)
    c.line.fill.background()

add_text(sl, "🗂",
         Inches(4.0), Inches(1.7), Inches(5.33), Inches(1.2),
         font_size=56, align=PP_ALIGN.CENTER)
add_text(sl, "WOW",
         Inches(2.0), Inches(2.8), Inches(9.33), Inches(1.0),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Work Overview Window",
         Inches(2.0), Inches(3.75), Inches(9.33), Inches(0.6),
         font_size=24, color=RGBColor(0xB3, 0xD4, 0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "담당자별 업무 진행 항목 / 현황 관리 · 격주 단위",
         Inches(2.0), Inches(4.35), Inches(9.33), Inches(0.45),
         font_size=16, color=RGBColor(0x99, 0xC2, 0xFF), align=PP_ALIGN.CENTER)

# 하단 기술 뱃지
for i, label in enumerate(["Microsoft Azure AD", "Firebase Firestore", "Microsoft Teams"]):
    x = Inches(2.8 + i * 2.7)
    add_rect(sl, x, Inches(5.5), Inches(2.4), Inches(0.4), RGBColor(0x00, 0x3E, 0x9E))
    add_text(sl, label, x, Inches(5.5), Inches(2.4), Inches(0.4),
             font_size=12, color=RGBColor(0xB3, 0xD4, 0xFF), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# 슬라이드 2 — 개요 (스크린샷: 상황판 펼침)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "개요", "WOW가 무엇인가요?")

# 왼쪽 텍스트
add_bullet_box(sl, [
    "담당자별 주간 일감을 격주 단위로 관리",
    "팀원의 상태(업무 중/휴가/종료)를 실시간 확인",
    "외부 관계자는 요약 화면에서 담당자 현황 파악",
    "Teams 업무 종료 보고 자동 생성 및 전송",
    "관리자 / 팀원 / 외부인 역할 기반 접근 제어",
], Inches(0.4), Inches(1.3), Inches(5.5), Inches(3.5), font_size=14)

add_text(sl, "Microsoft Azure AD 로그인 → 자동 멤버 등록\nFirebase Firestore 실시간 동기화",
         Inches(0.4), Inches(5.0), Inches(5.5), Inches(0.8),
         font_size=12, color=MUTED)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['status'], Inches(6.1), Inches(1.25), Inches(6.9), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 3 — 역할 구분
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "역할 구분", "관리자 · 팀원 · 외부인")

roles = [
    (AMBER,  "👑 관리자",  [
        "모든 기능 사용 가능",
        "담당자 추가 / 삭제 / 역할 변경",
        "타 담당자 일감 이동 및 복사",
        "외부인 뷰 미리보기",
        "담당자 순서 드래그 변경",
    ]),
    (BLUE,   "🙋 팀원",    [
        "본인 일감 및 상태 관리",
        "팀 전체 주간 보드 조회",
        "업무 설명 입력",
        "Teams 업무 종료 보고",
    ]),
    (PURPLE, "🌐 외부인",  [
        "팀 업무 요약(외부인 뷰)만 열람",
        "담당자 상태 및 업무 설명 확인",
        "태그 필터로 담당자 검색",
        "Teams로 담당자 직접 연락",
    ]),
]
for i, (color, title, items) in enumerate(roles):
    x = Inches(0.4 + i * 4.3)
    shape = sl.shapes.add_shape(1, x, Inches(1.3), Inches(4.0), Inches(5.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    add_rect(sl, x, Inches(1.3), Inches(4.0), Inches(0.55), color)
    add_text(sl, title, x + Inches(0.15), Inches(1.35), Inches(3.7), Inches(0.45),
             font_size=16, bold=True, color=WHITE)
    add_bullet_box(sl, items, x + Inches(0.15), Inches(2.0),
                   Inches(3.7), Inches(4.5), font_size=13)

add_text(sl, "※ 최초 가입 시 관리자가 이미 있으면 외부인으로 등록됩니다. 관리자가 역할을 변경해야 합니다.",
         Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
         font_size=11, color=MUTED)

# ════════════════════════════════════════════════
# 슬라이드 4 — 상황판 (스크린샷: 펼친 상황판)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "상황판", "팀 현황을 한눈에 확인")

# 왼쪽 텍스트
section_label(sl, "접기 / 펼치기", Inches(0.4), Inches(1.3))
add_bullet_box(sl, [
    "기본값: 접힌 상태",
    "헤더 클릭 → 펼치기 / 접기 토글",
    "접힌 상태: 업무 중 N명 / 휴가 N명 / 종료 N명 요약 표시",
], Inches(0.4), Inches(1.72), Inches(5.6), Inches(1.3), font_size=13)

section_label(sl, "내 상태 변경", Inches(0.4), Inches(3.1))
states = [
    (GREEN,                          "업무 중", "기본 상태"),
    (RGBColor(0x00, 0x8D, 0xA8),    "휴가",    "종료일 선택 → 당일 자정 이후 자동 복귀"),
    (MUTED,                          "종료 📤", "Teams 보고 후 상태 변경 / 다음 평일 자동 복귀"),
]
for i, (c, label, desc) in enumerate(states):
    y = Inches(3.55 + i * 0.82)
    add_rect(sl, Inches(0.4), y + Inches(0.08), Inches(0.15), Inches(0.15), c)
    add_text(sl, label, Inches(0.65), y, Inches(1.4), Inches(0.35),
             font_size=13, bold=True, color=c)
    add_text(sl, desc, Inches(2.15), y, Inches(3.9), Inches(0.35),
             font_size=12, color=DARK)

section_label(sl, "업무 설명 & 미보고 알림", Inches(0.4), Inches(5.9))
add_bullet_box(sl, [
    "본인 행 클릭 → 현재 업무 입력 (최대 60자)",
    "평일 18시 이후 업무 중 상태 → ⚠ 미보고 배지 표시",
], Inches(0.4), Inches(6.3), Inches(5.6), Inches(0.9), font_size=13)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['status'], Inches(6.2), Inches(1.25), Inches(6.8), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 5 — 주간 업무 보드 (스크린샷: 보드 접힌 상태)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "주간 업무 보드", "담당자별 일감 현황 관리")

section_label(sl, "보드 구성", Inches(0.4), Inches(1.3))
add_bullet_box(sl, [
    "담당자별 월~금 5개 컬럼 구성",
    "그룹 설정 시 그룹 헤더로 구분 표시",
    "주차 네비게이션으로 이전 / 다음 주 이동",
], Inches(0.4), Inches(1.72), Inches(5.6), Inches(1.3), font_size=13)

section_label(sl, "일감 추가", Inches(0.4), Inches(3.1))
add_bullet_box(sl, [
    "요일 컬럼 하단 ＋ 버튼 클릭",
    "여러 날 동시 추가 가능",
    "주 단위 반복 등록 가능",
], Inches(0.4), Inches(3.52), Inches(5.6), Inches(1.3), font_size=13)

section_label(sl, "일감 상태 (클릭으로 순환)", Inches(0.4), Inches(4.9))
for i, (c, label) in enumerate([
    (MUTED, "할 일"), (BLUE, "진행 중"), (GREEN, "완료"), (AMBER, "보류")
]):
    add_rect(sl, Inches(0.4 + i * 1.3), Inches(5.32), Inches(1.15), Inches(0.38), c)
    add_text(sl, label, Inches(0.4 + i * 1.3), Inches(5.32), Inches(1.15), Inches(0.38),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['board'], Inches(6.2), Inches(1.25), Inches(6.8), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 6 — 일감 이동 / 복사 (스크린샷: 드래그)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "일감 이동 & 복사", "드래그 앤 드롭으로 손쉽게")

section_label(sl, "일감 이동 (드래그 앤 드롭)", Inches(0.4), Inches(1.3))
add_bullet_box(sl, [
    "같은 담당자의 다른 요일 컬럼으로 드래그",
    "관리자는 다른 담당자 컬럼으로도 이동 가능",
    "드롭 위치에 파란 링 표시로 피드백 제공",
], Inches(0.4), Inches(1.72), Inches(5.6), Inches(1.4), font_size=13)

section_label(sl, "일감 복사 (관리자 전용)", Inches(0.4), Inches(3.2), AMBER)
add_bullet_box(sl, [
    "일감 우측 📋 버튼 클릭",
    "대상 담당자 / 주차 / 요일 선택 후 복사",
    "원본 일감은 그대로 유지됨",
], Inches(0.4), Inches(3.62), Inches(5.6), Inches(1.4), font_size=13)

section_label(sl, "이월 업무", Inches(0.4), Inches(5.1), RGBColor(0xE5, 0x6D, 0x00))
add_bullet_box(sl, [
    "완료되지 않은 업무를 다음 주로 이월",
    "주황색 배경 섹션으로 구분 표시",
    "드래그 앤 드롭으로 순서 변경 가능",
], Inches(0.4), Inches(5.52), Inches(5.6), Inches(1.4), font_size=13)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['drag'], Inches(6.2), Inches(1.25), Inches(6.8), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 7 — 담당자 관리 (스크린샷: 모달)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "담당자 관리", "관리자 전용 · 헤더 '👤 담당자 관리' 버튼")

section_label(sl, "담당자 추가 / 수정 항목", Inches(0.4), Inches(1.3))
fields = [
    ("이름",      "필수", BLUE),
    ("직급",      "필수", BLUE),
    ("이모지",    "프로필 아이콘 선택", MUTED),
    ("부서/그룹", "그룹 구분 표시에 사용됨", MUTED),
    ("태그",      "외부인 뷰 필터용 · 복수 입력 · 자동완성 지원", GREEN),
]
for i, (name, desc, c) in enumerate(fields):
    y = Inches(1.75 + i * 0.65)
    shape = sl.shapes.add_shape(1, Inches(0.4), y, Inches(5.6), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = c
    shape.line.width = Pt(1)
    add_text(sl, name, Inches(0.55), y + Inches(0.12), Inches(1.3), Inches(0.32),
             font_size=12, bold=True, color=c)
    add_text(sl, desc, Inches(2.0), y + Inches(0.12), Inches(3.8), Inches(0.32),
             font_size=12, color=DARK)

section_label(sl, "역할 변경 & 순서 변경", Inches(0.4), Inches(5.35))
add_bullet_box(sl, [
    "드롭다운으로 관리자 / 팀원 / 외부인 즉시 변경",
    "⠿ 핸들 드래그로 담당자 순서 변경",
    "마지막 관리자는 역할 변경 불가",
], Inches(0.4), Inches(5.78), Inches(5.6), Inches(1.3), font_size=13)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['manage'], Inches(6.2), Inches(1.25), Inches(6.8), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 8 — Teams 연동
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "Teams 연동", "업무 종료 보고 · 외부인 Teams 연락")

# 종료 보고 카드
shape = sl.shapes.add_shape(1, Inches(0.4), Inches(1.3), Inches(5.9), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = BLUE
shape.line.width = Pt(1.5)

add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.5), BLUE)
add_text(sl, "📤  업무 종료 보고 흐름",
         Inches(0.6), Inches(1.33), Inches(5.5), Inches(0.44),
         font_size=15, bold=True, color=WHITE)

steps = [
    ("1", "상황판에서 '종료 📤' 버튼 클릭"),
    ("2", "오늘 완료 업무 목록 자동 채워진 양식 열림"),
    ("3", "내용 자유 수정 가능"),
    ("4", "'Teams 전송 후 종료' 클릭"),
    ("5", "전송 완료 → 상태 '종료'로 변경"),
]
for i, (num, text) in enumerate(steps):
    y = Inches(2.0 + i * 0.82)
    add_rect(sl, Inches(0.6), y + Inches(0.06), Inches(0.3), Inches(0.3), BLUE)
    add_text(sl, num, Inches(0.6), y + Inches(0.04), Inches(0.3), Inches(0.3),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, text, Inches(1.05), y, Inches(4.9), Inches(0.45),
             font_size=13, color=DARK)

add_text(sl, "※ '종료만 처리': 전송 없이 상태만 변경\n※ '취소': 창 닫기, 업무 중 상태 유지",
         Inches(0.6), Inches(6.3), Inches(5.5), Inches(0.45),
         font_size=11, color=MUTED)

# 외부인 Teams 연락 카드
shape2 = sl.shapes.add_shape(1, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.5))
shape2.fill.solid()
shape2.fill.fore_color.rgb = WHITE
shape2.line.color.rgb = PURPLE
shape2.line.width = Pt(1.5)

add_rect(sl, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.5), PURPLE)
add_text(sl, "💬  외부인 Teams 연락",
         Inches(7.0), Inches(1.33), Inches(5.7), Inches(0.44),
         font_size=15, bold=True, color=WHITE)
add_bullet_box(sl, [
    "외부인 뷰 멤버 카드 하단에 메시지 입력창 표시",
    "메시지 입력 후 💬 버튼 또는 Enter 클릭",
    "Teams 앱이 열리며 채팅창에 메시지 입력된 상태로 시작",
    "이메일 등록된 담당자만 버튼 표시",
    "종료 / 휴가 상태 담당자는 버튼 미표시",
], Inches(7.0), Inches(2.0), Inches(5.7), Inches(4.5), font_size=13)

# ════════════════════════════════════════════════
# 슬라이드 9 — 외부인 뷰 (스크린샷)
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BG)
header_bar(sl, "외부인 뷰", "외부인 전용 팀 업무 요약 화면")

section_label(sl, "화면 구성", Inches(0.4), Inches(1.3))
add_bullet_box(sl, [
    "상단 배너: 현재 업무 중 인원 수 표시",
    "멤버 카드 그리드 (최대 4열)",
    "카드: 이모지 · 상태 · 이름 · 직급 · 태그 · 업무 설명",
    "상태 순서: 업무 중 → 휴가 → 종료",
], Inches(0.4), Inches(1.72), Inches(5.6), Inches(1.8), font_size=13)

section_label(sl, "태그 필터", Inches(0.4), Inches(3.6))
add_bullet_box(sl, [
    "뷰 상단 태그 버튼 클릭 → 해당 태그 담당자만 표시",
    "여러 태그 동시 선택 (AND 조건)",
    "'✕ 초기화' 버튼으로 필터 해제",
], Inches(0.4), Inches(4.02), Inches(5.6), Inches(1.3), font_size=13)

section_label(sl, "관리자 미리보기", Inches(0.4), Inches(5.4))
add_text(sl, "헤더 '👁 외부인 뷰' 버튼으로 관리자도 외부인 화면 확인 가능",
         Inches(0.6), Inches(5.78), Inches(5.6), Inches(0.5),
         font_size=13, color=DARK)

# 오른쪽 스크린샷
add_screenshot(sl, IMG['external'], Inches(6.2), Inches(1.25), Inches(6.8), Inches(5.9))

# ════════════════════════════════════════════════
# 슬라이드 10 — 마무리
# ════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl, BLUE)

for (x, y, s) in [(Inches(10), Inches(0.3), Inches(4.5)),
                   (Inches(-0.5), Inches(4.5), Inches(4))]:
    c = sl.shapes.add_shape(9, x, y, s, s)
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(0x00, 0x3E, 0x9E)
    c.line.fill.background()

add_text(sl, "🗂",
         Inches(3.5), Inches(1.8), Inches(6.33), Inches(1.1),
         font_size=52, align=PP_ALIGN.CENTER)
add_text(sl, "WOW",
         Inches(2.0), Inches(2.85), Inches(9.33), Inches(1.0),
         font_size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Work Overview Window",
         Inches(2.0), Inches(3.8), Inches(9.33), Inches(0.6),
         font_size=22, color=RGBColor(0xB3, 0xD4, 0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "담당자별 업무 진행 항목 / 현황 관리 · 격주 단위",
         Inches(2.0), Inches(4.45), Inches(9.33), Inches(0.4),
         font_size=15, color=RGBColor(0x99, 0xC2, 0xFF), align=PP_ALIGN.CENTER)

out = r"C:\Users\sam.lee\Desktop\잡일\WOW\WOW\WOW_Manual.pptx"
prs.save(out)
print("저장 완료:", out)
